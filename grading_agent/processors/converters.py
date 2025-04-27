# -*- coding: utf-8 -*-
"""Asynchronous file content converters for various formats."""

import asyncio
import base64
import io
import logging
import os
import tempfile
from io import BytesIO
from typing import List, Optional, Tuple

import nbformat
from docx import Document
from nbconvert import PDFExporter
from pdf2image import convert_from_path
from PIL import Image, ImageDraw, UnidentifiedImageError
from pptx import Presentation

# Import aiofiles for async file operations
try:
    import aiofiles
except ImportError:
    aiofiles = None # type: ignore # Handle optional dependency

logger = logging.getLogger(__name__)

# --- Helper Functions (Synchronous, potentially blocking) ---

def _compress_image_sync(image: Image.Image, max_size_mb: int = 2) -> BytesIO:
    """Synchronously compresses a PIL Image to JPEG format below a size limit.

    Args:
        image: The PIL Image object.
        max_size_mb: The maximum desired size in megabytes.

    Returns:
        A BytesIO object containing the compressed JPEG image data.
    """
    buffered = io.BytesIO()
    # Convert to RGB if necessary (JPEG doesn't support alpha)
    if image.mode == 'RGBA':
        try:
            # Create a white background image
            bg = Image.new("RGB", image.size, (255, 255, 255))
            # Paste the RGBA image onto the white background
            bg.paste(image, mask=image.split()[3]) # 3 is the alpha channel
            image = bg
            logger.debug("Converted RGBA image to RGB with white background.")
        except Exception as e:
            logger.warning("Failed to convert RGBA to RGB, using direct conversion: %s", e)
            image = image.convert('RGB') # Fallback

    elif image.mode != 'RGB':
         image = image.convert('RGB')


    # Initial save with high quality
    try:
        image.save(buffered, format="JPEG", quality=95, optimize=True)
        size_kb = len(buffered.getvalue()) / 1024
        logger.debug("Initial image size: %.2f KB", size_kb)

        # Reduce quality if size exceeds limit
        if size_kb > max_size_mb * 1024:
            quality = 95
            low, high = 10, 95 # Range for quality search
            best_quality_buffer = buffered # Keep the last successful buffer

            # Iteratively reduce quality - more robust than direct calculation
            while high >= low and size_kb > max_size_mb * 1024:
                quality = (low + high) // 2
                if quality <= 10: # Don't go too low
                    logger.warning("Image compression hit minimum quality (10), size might still exceed limit.")
                    break

                temp_buffer = io.BytesIO()
                image.save(temp_buffer, format="JPEG", quality=quality, optimize=True)

                if len(temp_buffer.getvalue()) <= max_size_mb * 1024:
                    best_quality_buffer = temp_buffer # Found a valid size
                    low = quality + 1 # Try slightly higher quality
                else:
                    high = quality - 1 # Need lower quality

                size_kb = len(best_quality_buffer.getvalue()) / 1024
                buffered = best_quality_buffer # Use the best buffer found so far

            logger.info("Compressed image to %.2f KB with quality %d", size_kb, quality)

    except Exception as e:
        logger.error("Error during image compression: %s", e)
        # Return an empty buffer or raise? Returning empty for now.
        return io.BytesIO()

    buffered.seek(0) # Reset buffer position
    return buffered


def _image_to_base64_sync(image: Image.Image, max_size_mb: int = 2) -> str:
    """Synchronously converts a PIL Image to a Base64 encoded string.

    Args:
        image: The PIL Image object.
        max_size_mb: Maximum size for compression.

    Returns:
        A Base64 encoded string of the compressed image.
    """
    try:
        buffered = _compress_image_sync(image, max_size_mb)
        return base64.b64encode(buffered.getvalue()).decode('utf-8')
    except Exception as e:
        logger.error("Error encoding image to base64: %s", e)
        return ""


def _convert_ipynb_to_pdf_sync(ipynb_path: str) -> Optional[str]:
    """Synchronously converts an ipynb file to a temporary PDF file path.

    Args:
        ipynb_path: Path to the .ipynb file.

    Returns:
        Path to the created temporary PDF file, or None on failure.
    """
    try:
        nb = None
        # Try common encodings
        for encoding in ['utf-8', 'latin1', 'cp1252', 'gbk']:
            try:
                with open(ipynb_path, 'r', encoding=encoding) as f:
                    nb = nbformat.read(f, as_version=4)
                logger.debug("Successfully read ipynb %s with encoding %s", ipynb_path, encoding)
                break
            except UnicodeDecodeError:
                continue # Try next encoding
            except FileNotFoundError:
                 logger.error("ipynb file not found: %s", ipynb_path)
                 return None
            except Exception as e:
                # Log other read errors but continue trying encodings
                logger.warning("Error reading ipynb %s with encoding %s: %s", ipynb_path, encoding, e)
                continue

        if nb is None:
             logger.error("Failed to read ipynb %s with any attempted encoding.", ipynb_path)
             return None

        # Configure PDF export
        pdf_exporter = PDFExporter()
        # pdf_exporter.template_name = 'classic' # Or specify a custom template if needed
        # Consider adding configuration for template if necessary

        try:
            (pdf_data, _) = pdf_exporter.from_notebook_node(nb)
            # Create a temporary file that persists until explicitly deleted
            # Use NamedTemporaryFile from tempfile module
            temp_pdf_file = tempfile.NamedTemporaryFile(delete=False, suffix=".pdf")
            temp_pdf_file.write(pdf_data)
            temp_pdf_file.close() # Close the file handle
            pdf_path = temp_pdf_file.name
            logger.info("Successfully converted %s to temporary PDF: %s", ipynb_path, pdf_path)
            return pdf_path
        except Exception as e:
            logger.error("Error converting notebook node to PDF for %s: %s", ipynb_path, e)
            # Clean up temp file if conversion failed mid-way
            if 'temp_pdf_file' in locals() and os.path.exists(temp_pdf_file.name):
                 os.remove(temp_pdf_file.name)
            return None

    except Exception as e:
        # Catch-all for unexpected errors during the process
        logger.error("Failed to process ipynb file %s: %s", ipynb_path, e)
        return None


def _convert_pdf_to_base64_images_sync(pdf_path: str) -> List[str]:
    """Synchronously converts a PDF file to a list of Base64 encoded images.

    Args:
        pdf_path: Path to the PDF file.

    Returns:
        A list of Base64 encoded strings, one for each page image.
    """
    base64_images: List[str] = []
    if not os.path.exists(pdf_path):
        logger.error("PDF file not found: %s", pdf_path)
        return base64_images
    try:
        # Use poppler_path if specified in environment or config
        poppler_path = os.environ.get("POPPLER_PATH")
        images = convert_from_path(pdf_path, poppler_path=poppler_path)

        if not images:
            logger.warning("No pages converted from PDF: %s", pdf_path)
            return base64_images

        logger.info("Converting %d pages from PDF: %s", len(images), pdf_path)
        for i, image in enumerate(images):
            try:
                img_base64 = _image_to_base64_sync(image)
                if img_base64:
                    base64_images.append(img_base64)
            except Exception as e:
                logger.error("Error processing page %d from PDF %s: %s", i + 1, pdf_path, e)
        return base64_images
    except Exception as e:
        # Catch errors from convert_from_path (e.g., Poppler not found, corrupted PDF)
        logger.error("Failed to convert PDF %s to images: %s", pdf_path, e)
        return base64_images # Return empty list on failure


def _convert_docx_to_content_sync(docx_path: str) -> Tuple[List[str], List[str]]:
    """Synchronously extracts text and renders content as images from DOCX.

    Args:
        docx_path: Path to the .docx file.

    Returns:
        A tuple containing:
        - List of text content strings (paragraphs).
        - List of Base64 encoded image strings (rendered content, embedded images).
    """
    texts: List[str] = []
    base64_images: List[str] = []
    try:
        doc = Document(docx_path)

        # 1. Extract Text
        for para in doc.paragraphs:
            if para.text.strip():
                texts.append(para.text.strip())
        logger.info("Extracted %d paragraphs from DOCX: %s", len(texts), docx_path)

        # 2. Extract Embedded Images (Simplified)
        # python-docx doesn't make image extraction straightforward across all rels.
        # This is a basic attempt; more robust extraction might need xml parsing.
        image_parts = [part for rel in doc.part.rels.values()
                       if "image" in rel.target_ref.lower()
                       for part in [rel.target_part] if hasattr(part, 'blob')]

        logger.info("Found %d potential image parts in DOCX: %s", len(image_parts), docx_path)
        for i, img_part in enumerate(image_parts):
             try:
                 img_bytes = img_part.blob
                 img = Image.open(BytesIO(img_bytes))
                 img_base64 = _image_to_base64_sync(img)
                 if img_base64:
                     base64_images.append(img_base64)
                 logger.debug("Extracted embedded image %d from DOCX: %s", i+1, docx_path)
             except UnidentifiedImageError:
                 logger.warning("Could not identify image format for part %d in DOCX: %s", i+1, docx_path)
             except Exception as e:
                 logger.error("Error processing embedded image %d from DOCX %s: %s", i+1, docx_path, e)


        # 3. Render Text as Image (Fallback/Optional - based on original prepare.py)
        # This might be redundant if text extraction is sufficient. Kept for compatibility.
        # Consider making this optional via configuration.
        render_as_image = False # Set to True to enable rendering like prepare.py
        if render_as_image and texts:
            logger.info("Rendering DOCX text content as image for: %s", docx_path)
            height_per_paragraph = 20 # Smaller height for rendering
            min_height = 200
            max_height = 10000 # Limit max image height
            width = 800
            total_height = min(max_height, max(min_height, len(texts) * height_per_paragraph))

            img = Image.new('RGB', (width, total_height), color='white')
            draw = ImageDraw.Draw(img)
            current_y = 10

            for text_line in texts:
                if current_y >= total_height - height_per_paragraph:
                    logger.warning("DOCX rendering reached max height for %s", docx_path)
                    break
                try:
                    # Basic text drawing
                    draw.text((10, current_y), text_line, fill='black')
                    current_y += height_per_paragraph
                except Exception as e:
                     logger.warning("Error rendering text line in %s: %s", docx_path, e)

            if current_y > 10: # Only save if content was drawn
                try:
                    rendered_img_base64 = _image_to_base64_sync(img)
                    if rendered_img_base64:
                        base64_images.insert(0, rendered_img_base64) # Add rendered image first
                    logger.info("Successfully rendered DOCX text %s to image.", docx_path)
                except Exception as e:
                    logger.error("Error converting rendered DOCX image to base64 for %s: %s", docx_path, e)

        return texts, base64_images
    except Exception as e:
        logger.error("Failed to process DOCX file %s: %s", docx_path, e)
        return texts, base64_images # Return whatever was extracted


def _convert_pptx_to_content_sync(pptx_path: str) -> Tuple[List[str], List[str]]:
    """Synchronously extracts text and images from PPTX slides.

    Args:
        pptx_path: Path to the .pptx file.

    Returns:
        A tuple containing:
        - List of text content strings (from text shapes).
        - List of Base64 encoded image strings (embedded images).
    """
    texts: List[str] = []
    base64_images: List[str] = []
    try:
        prs = Presentation(pptx_path)
        logger.info("Processing %d slides from PPTX: %s", len(prs.slides), pptx_path)

        for i, slide in enumerate(prs.slides):
            slide_texts: List[str] = []
            logger.debug("Processing slide %d", i + 1)
            for shape in slide.shapes:
                 try:
                    # Extract Text
                    if shape.has_text_frame and shape.text.strip():
                        text = shape.text.strip()
                        slide_texts.append(text)
                        logger.debug("  Extracted text: %.50s...", text)

                    # Extract Images
                    if hasattr(shape, 'image'):
                        try:
                            image_stream = BytesIO(shape.image.blob)
                            img = Image.open(image_stream)
                            img_base64 = _image_to_base64_sync(img)
                            if img_base64:
                                base64_images.append(img_base64)
                            logger.debug("  Extracted image shape (%.2f KB)", len(img_base64)*3/4/1024)
                        except UnidentifiedImageError:
                            logger.warning("  Could not identify image format for shape on slide %d", i+1)
                        except Exception as img_e:
                            logger.warning("  Error processing image shape on slide %d: %s", i+1, img_e)

                 except Exception as shape_e:
                     logger.warning("  Error processing shape on slide %d: %s", i+1, shape_e)

            if slide_texts:
                 texts.append(f"--- Slide {i+1} ---\n" + "\n".join(slide_texts))

        logger.info("Extracted %d text blocks and %d images from PPTX: %s", len(texts), len(base64_images), pptx_path)
        return texts, base64_images
    except Exception as e:
        logger.error("Failed to process PPTX file %s: %s", pptx_path, e)
        return texts, base64_images # Return whatever was extracted


# --- Async Wrappers ---

async def compress_image(image: Image.Image, max_size_mb: int = 2) -> BytesIO:
    """Asynchronously compresses a PIL Image."""
    return await asyncio.to_thread(_compress_image_sync, image, max_size_mb)

async def image_to_base64(image: Image.Image, max_size_mb: int = 2) -> str:
    """Asynchronously converts a PIL Image to a Base64 string."""
    return await asyncio.to_thread(_image_to_base64_sync, image, max_size_mb)

async def convert_ipynb_to_pdf(ipynb_path: str) -> Optional[str]:
    """Asynchronously converts an ipynb file to a temporary PDF file path."""
    logger.info("Starting async conversion of ipynb to PDF: %s", ipynb_path)
    pdf_path = await asyncio.to_thread(_convert_ipynb_to_pdf_sync, ipynb_path)
    if pdf_path:
        logger.info("Finished async conversion of ipynb to PDF: %s -> %s", ipynb_path, pdf_path)
    else:
        logger.error("Async conversion of ipynb to PDF failed for: %s", ipynb_path)
    return pdf_path


async def convert_pdf_to_base64_images(pdf_path: str) -> List[str]:
    """Asynchronously converts a PDF file to a list of Base64 encoded images."""
    logger.info("Starting async conversion of PDF to images: %s", pdf_path)
    result = await asyncio.to_thread(_convert_pdf_to_base64_images_sync, pdf_path)
    logger.info("Finished async conversion of PDF %s. Got %d images.", pdf_path, len(result))
    return result

async def convert_docx_to_content(docx_path: str) -> Tuple[List[str], List[str]]:
    """Asynchronously extracts content (text, images) from a DOCX file."""
    logger.info("Starting async content extraction from DOCX: %s", docx_path)
    texts, images = await asyncio.to_thread(_convert_docx_to_content_sync, docx_path)
    logger.info("Finished async extraction from DOCX %s. Got %d texts, %d images.", docx_path, len(texts), len(images))
    return texts, images

async def convert_pptx_to_content(pptx_path: str) -> Tuple[List[str], List[str]]:
    """Asynchronously extracts content (text, images) from a PPTX file."""
    logger.info("Starting async content extraction from PPTX: %s", pptx_path)
    texts, images = await asyncio.to_thread(_convert_pptx_to_content_sync, pptx_path)
    logger.info("Finished async extraction from PPTX %s. Got %d text blocks, %d images.", pptx_path, len(texts), len(images))
    return texts, images


# --- Main Conversion Function ---

async def convert_file_content(
    file_path: str, mime_type: str
) -> Tuple[List[str], List[str]]:
    """
    Detects file type and converts content asynchronously.

    Args:
        file_path: Path to the file.
        mime_type: MIME type of the file.

    Returns:
        A tuple containing:
        - List of text content strings.
        - List of Base64 encoded image strings.
    """
    texts: List[str] = []
    images: List[str] = []
    temp_pdf_path: Optional[str] = None # To track temporary PDF for cleanup

    try:
        logger.info("Converting file: %s with MIME type: %s", file_path, mime_type)

        if not os.path.exists(file_path):
             logger.error("File not found for conversion: %s", file_path)
             return texts, images

        if mime_type == 'application/x-ipynb+json' or file_path.endswith('.ipynb'):
            temp_pdf_path = await convert_ipynb_to_pdf(file_path)
            if temp_pdf_path:
                # Convert the generated PDF to images (text extraction from PDF not implemented here)
                images = await convert_pdf_to_base64_images(temp_pdf_path)
            else:
                 logger.error("Failed to convert ipynb to PDF, cannot extract content: %s", file_path)

        elif mime_type == 'application/pdf':
            # Only extract images for now. Text extraction could be added.
            images = await convert_pdf_to_base64_images(file_path)
            # Example: texts = await extract_text_from_pdf(file_path) # Needs implementation

        elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
            texts, images = await convert_docx_to_content(file_path)

        elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
            texts, images = await convert_pptx_to_content(file_path)

        elif mime_type.startswith('image/'):
            # Use asyncio.to_thread for potentially blocking file read and image processing
            def read_and_encode_image_sync():
                try:
                    with Image.open(file_path) as img:
                        return _image_to_base64_sync(img)
                except UnidentifiedImageError:
                     logger.error("Cannot identify image file: %s", file_path)
                     return None
                except FileNotFoundError:
                     logger.error("Image file not found during sync read: %s", file_path)
                     return None
                except Exception as e:
                     logger.error("Error processing image file %s in sync thread: %s", file_path, e)
                     return None

            img_base64 = await asyncio.to_thread(read_and_encode_image_sync)
            if img_base64:
                images.append(img_base64)

        elif mime_type.startswith('text/') or \
             mime_type in ('application/octet-stream', 'application/x-python', 'application/javascript'): # Add common code types
             # Use aiofiles for async text reading if available
             if aiofiles:
                 try:
                     async with aiofiles.open(file_path, mode='r', encoding='utf-8', errors='ignore') as f:
                         content = await f.read()
                         texts.append(content)
                     logger.info("Read text content from: %s", file_path)
                 except FileNotFoundError:
                     logger.error("Text file not found during async read: %s", file_path)
                 except Exception as e:
                     logger.error("Error reading text file %s with aiofiles: %s", file_path, e)
             else:
                 # Fallback to sync reading in thread if aiofiles not installed
                 logger.warning("aiofiles not installed, using sync file reading in thread for %s.", file_path)
                 def read_text_sync():
                     try:
                         with open(file_path, mode='r', encoding='utf-8', errors='ignore') as f:
                             return f.read()
                     except FileNotFoundError:
                         logger.error("Text file not found during sync read: %s", file_path)
                         return None
                     except Exception as e:
                         logger.error("Error reading text file %s in sync thread: %s", file_path, e)
                         return None
                 content = await asyncio.to_thread(read_text_sync)
                 if content is not None:
                      texts.append(content)
        else:
            logger.warning("Unsupported MIME type for conversion: %s (%s)", mime_type, file_path)

    finally:
        # Clean up temporary PDF file if created
        if temp_pdf_path and os.path.exists(temp_pdf_path):
            try:
                # Run cleanup in a separate thread to avoid blocking
                await asyncio.to_thread(os.remove, temp_pdf_path)
                logger.debug("Cleaned up temporary PDF: %s", temp_pdf_path)
            except Exception as e:
                logger.error("Failed to clean up temporary PDF %s: %s", temp_pdf_path, e)

    logger.info("Conversion result for %s: %d text blocks, %d images.", file_path, len(texts), len(images))
    return texts, images