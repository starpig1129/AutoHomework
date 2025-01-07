import nbformat
import io
import os
import tempfile
import base64
from nbconvert import PDFExporter
from PIL import Image, ImageDraw
from docx import Document
from pptx import Presentation
from io import BytesIO
from pdf2image import convert_from_path
def compress_image(image, max_size_mb=2):
    """Compress image to a maximum size of 2MB for API compatibility."""
    buffered = io.BytesIO()
    # Convert to RGB if image is in RGBA mode
    if image.mode == 'RGBA':
        image = image.convert('RGB')
    
    # Initial save with high quality
    image.save(buffered, format="JPEG", quality=95)
    size_kb = len(buffered.getvalue()) / 1024
    
    if size_kb > max_size_mb * 1024:
        # Calculate compression ratio
        compression_ratio = (max_size_mb * 1024) / size_kb
        quality = int(95 * compression_ratio)
        
        # Resize if quality would be too low
        if quality < 30:
            scale_factor = (max_size_mb * 1024 * 30) / (size_kb * 95)
            new_width = int(image.width * scale_factor)
            new_height = int(image.height * scale_factor)
            image = image.resize((new_width, new_height), Image.LANCZOS)
            quality = 30
        
        # Save with new quality/size
        buffered = io.BytesIO()
        image.save(buffered, format="JPEG", quality=quality)
        
    return buffered
# 定義函數：將 ipynb 轉換為 PDF
def convert_ipynb_to_pdf(ipynb_path):
    try:
        # Try different encodings
        for encoding in ['utf-8', 'latin1', 'cp1252']:
            try:
                with open(ipynb_path, 'r', encoding=encoding) as f:
                    nb = nbformat.read(f, as_version=4)
                break
            except UnicodeDecodeError:
                continue
            except Exception as e:
                print(f"讀取 notebook 時發生錯誤 (編碼: {encoding}): {str(e)}")
                continue
        else:
            raise ValueError("無法以任何編碼讀取檔案")

        pdf_exporter = PDFExporter()
        pdf_exporter.template_name = 'classic'
        
        try:
            pdf_data, _ = pdf_exporter.from_notebook_node(nb)
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pdf') as tf:
                tf.write(pdf_data)
                pdf_path = tf.name
            return pdf_path
        except Exception as e:
            print(f"轉換 PDF 時發生錯誤: {str(e)}")
            return None
            
    except Exception as e:
        print(f"無法讀取的檔案: {ipynb_path}, 錯誤類型: {type(e).__name__}, 錯誤: {str(e)}")
        return None


# 將 PDF 轉換為 Base64 圖片
def convert_pdf_to_base64_images(pdf_path):
    try:
        # Check if file exists and is readable
        if not os.path.exists(pdf_path):
            print(f"PDF 檔案不存在: {pdf_path}")
            return None
            
        try:
            images = convert_from_path(pdf_path)
        except Exception as e:
            print(f"PDF 轉換圖片失敗: {str(e)}")
            return None
            
        if not images:
            print(f"PDF 沒有可轉換的頁面: {pdf_path}")
            return None
            
        base64_images = []
        for page_num, image in enumerate(images, 1):
            try:
                buffered = compress_image(image)
                img_base64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
                base64_images.append(img_base64)
            except Exception as e:
                print(f"處理 PDF 第 {page_num} 頁時發生錯誤: {str(e)}")
                continue
                
        return base64_images if base64_images else None
        
    except Exception as e:
        print(f"無法轉換的檔案: {pdf_path}, 錯誤類型: {type(e).__name__}, 錯誤: {str(e)}")
        return None


# 將 .docx 轉換為 Base64 圖片
def convert_docx_to_base64_images(docx_path):
    try:
        doc = Document(docx_path)
        base64_images = []
        
        # Calculate total content height
        total_paragraphs = len(doc.paragraphs)
        height_per_paragraph = 50  # Adjust based on content
        total_height = max(600, total_paragraphs * height_per_paragraph)
        
        # Create single image for all content
        img = Image.new('RGB', (800, total_height), color='white')
        d = ImageDraw.Draw(img)
        
        current_y = 10
        for paragraph in doc.paragraphs:
            try:
                # Handle potential encoding issues
                text = paragraph.text.encode('utf-8', errors='ignore').decode('utf-8')
                if text.strip():  # Only process non-empty paragraphs
                    d.text((10, current_y), text, fill='black')
                    current_y += height_per_paragraph
            except Exception as e:
                print(f"段落處理錯誤: {str(e)}")
                continue
        
        # Only create image if there's content
        if current_y > 10:
            buffered = compress_image(img)
            img_base64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
            base64_images.append(img_base64)
            
        return base64_images
    except Exception as e:
        print(f"無法讀取檔案: {docx_path}, 錯誤類型: {type(e).__name__}, 錯誤: {str(e)}")


# 將 .pptx 轉換為 Base64 圖片
def convert_pptx_to_base64_images(pptx_path):
    try:
        prs = Presentation(pptx_path)
        base64_images = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            try:
                # Create a new image with white background
                img = Image.new('RGB', (1920, 1080), color='white')
                d = ImageDraw.Draw(img)
                
                # Track position for content placement
                y_position = 10
                
                for shape in slide.shapes:
                    try:
                        # Handle images
                        if hasattr(shape, 'image'):
                            try:
                                # Extract image data
                                image_stream = BytesIO(shape.image.blob)
                                shape_img = Image.open(image_stream)
                                
                                # Calculate position to paste image
                                x = shape.left if hasattr(shape, 'left') else 10
                                y = shape.top if hasattr(shape, 'top') else y_position
                                
                                # Resize if needed while maintaining aspect ratio
                                if shape.width and shape.height:
                                    shape_img = shape_img.resize(
                                        (min(int(shape.width), 1900), 
                                         min(int(shape.height), 1000)), 
                                        Image.LANCZOS
                                    )
                                
                                # Paste image onto slide
                                img.paste(shape_img, (int(x), int(y)))
                                y_position = max(y_position, y + shape_img.height + 10)
                                
                            except Exception as e:
                                print(f"圖片處理錯誤: {str(e)}")
                                continue
                        
                        # Handle text
                        if hasattr(shape, "text"):
                            try:
                                text = shape.text.encode('utf-8', errors='ignore').decode('utf-8')
                                if text.strip():
                                    # Get text position if available
                                    x = shape.left if hasattr(shape, 'left') else 10
                                    y = shape.top if hasattr(shape, 'top') else y_position
                                    
                                    d.text((int(x), int(y)), text, fill='black')
                                    y_position = max(y_position, y + 30)
                            except Exception as e:
                                print(f"文字處理錯誤: {str(e)}")
                                continue
                                
                    except Exception as e:
                        print(f"形狀處理錯誤: {str(e)}")
                        continue
                
                # Convert slide to base64
                buffered = compress_image(img)
                img_base64 = base64.b64encode(buffered.getvalue()).decode('utf-8')
                base64_images.append(img_base64)
                
            except Exception as e:
                print(f"投影片 {slide_num} 處理錯誤: {str(e)}")
                continue
                
        return base64_images if base64_images else None
    except Exception as e:
        print(f"無法讀取的檔案: {pptx_path}, 錯誤類型: {type(e).__name__}, 錯誤: {str(e)}")


# 將圖片轉換為 Base64 字符串
def image_to_base64(image):
    buffered = compress_image(image)
    return base64.b64encode(buffered.getvalue()).decode('utf-8')
